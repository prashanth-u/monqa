from flask import Flask, request, abort, send_file
from flask_cors import CORS
import nltk
import json
import math
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk.corpus import stopwords
from nltk.stem import PorterStemmer
import glob
import fitz
import datetime
import os
import uuid
from pptx import Presentation
from werkzeug.utils import secure_filename
from flask_pymongo import pymongo

nltk.download('stopwords')
nltk.download('punkt')
CONNECTION_STRING = "mongodb+srv://me:me@monqa-nfobo.mongodb.net/test?retryWrites=true&w=majority"
client = pymongo.MongoClient(CONNECTION_STRING)
db = client.get_database('test')


stop = set(stopwords.words('english'))
stop.update(['.', ',', '"', "'", '?', '!', ':', ';', '(', ')', '[', ']', '{', '}']) # remove it if you need punctuation
ps = PorterStemmer()

app = Flask(__name__)

cors = CORS(app)


ALLOWED_EXTENSIONS = {'pdf', 'ppt', 'pptx'}

def allowed_file(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


def normalizedVectorValues(dv,qv):
    nDoc = 0.0
    nQuery = 0.0
    for val in dv:
        nDoc = nDoc + math.pow(dv[val],2)
    for val1 in qv:
        nQuery = nQuery + math.pow(qv[val1],2)
    return math.sqrt(nDoc) * math.sqrt(nQuery);

def removeStopWordsAndStemQuery(words):
    listOfTokens =  [word for word in words if word not in stop]
    stemmedTokens = []
    for w in listOfTokens:
        stemmedWord = ps.stem(w)
        stemmedTokens.append(stemmedWord)
    tf = getTf(stemmedTokens)
    return [stemmedTokens, tf]

def getTf(words):
    wordfreq = {}
    visited = []
    for w in words:
        if w not in visited:
            visited.append(w)
            wordfreq[w] = words.count(w)
    return wordfreq

def removeStopWordsAndStem(words, dictionary):
    listOfTokens =  [word for word in words if word not in stop]
    stemmedTokens = []
    for w in listOfTokens:
        stemmedWord = ps.stem(w)
        stemmedTokens.append(stemmedWord)
        if stemmedWord not in dictionary:
            dictionary.append(stemmedWord)
    tf = getTf(stemmedTokens)
    return [stemmedTokens, tf, dictionary]

def indexDocs(unit):
    indexId = uuid.uuid4()
    db.index.insert({'unit': unit, 'timestamp':datetime.datetime.now(), 'status': 'In Progress', 'uid': indexId})
    documents={}
    dictionary = []
    idf = {}
    pdfFileLocations = glob.glob("docsearch/files/"+unit+"/PDF/*.pdf")
    pptFileLocations = glob.glob("docsearch/files/"+unit+"/ppt/*.pptx")
    print(pptFileLocations)
    for j in pdfFileLocations:
        try:
            doc = fitz.Document(j)
        except Exception as e:
            print(e)
        nameOfDoc = j.split("/")[4]
        for i in range (0, doc.pageCount):
            page = doc.loadPage(i)
            pageText = page.getText("text")
            details = removeStopWordsAndStem(word_tokenize(pageText.encode('utf-16', 'surrogatepass').decode('utf-16')), dictionary)
            documents[nameOfDoc + "::::" + str(i+1)] = {"location": j, "documentName": nameOfDoc, "pageNo": i+1, "id": nameOfDoc + str(i+1), "text": pageText,"tokens": details[0], "tf": details[1]}

    for k in pptFileLocations:
        prs = Presentation(k)
        nameOfDoc = k.split("/")[4]
        for slide in prs.slides:
            pageText = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    print(prs.slides.index(slide)+1)
                    print(shape.text)
                    pageText = pageText + shape.text
            details = removeStopWordsAndStem(word_tokenize(pageText.encode('utf-16', 'surrogatepass').decode('utf-16')), dictionary)
            documents[nameOfDoc + "::::" + str(prs.slides.index(slide)+1)] = {"location": k, "documentName": nameOfDoc, "pageNo": prs.slides.index(slide)+1, "id": nameOfDoc + str(prs.slides.index(slide)+1), "text": pageText,
                                           "tokens": details[0], "tf": details[1]}
    indexFile = ""
    for word in dictionary:
        indexFile = indexFile + word
        count = 0
        for key in documents:
            uniqueWords = documents[key]['tf']
            if word in uniqueWords.keys():
                indexFile = indexFile + "???" + key + "???" + str(uniqueWords[word])
                count = count+1
        idf[word] = math.log(len(documents)/(count+1))
        indexFile = indexFile + "???" + str(idf[word]) + "\n"

    file = open("docsearch/files/"+unit+"/index.txt", "w")

    file.write(indexFile)
    file.close()

    indexRecord = db.index.find_one({'uid': indexId})

    if indexRecord is not None:
        indexRecord['status'] = "Success"
        db.index.save(indexRecord)
    return("success")


@app.route('/index/<unit>/<otp>', methods=['POST'])
def index(unit, otp):
    otpDocument = db.otps.find_one({'otpString': otp, 'used': False, 'role': 'Admin'})
    if otpDocument!= None and len(otpDocument) > 0:
        otpDocument['used'] = True
        db.otps.update_one({'_id':otpDocument['_id']}, {"$set": otpDocument}, upsert=False)
        return indexDocs(unit)
    else:
        return '',403



@app.route('/query/<unit>/<otp>', methods=['POST'])
def processQuery(unit,otp):
    
    query = request.json['question']
    otpDocument = db.otps.find_one({'otpString': otp, 'used': False}) 
    if otpDocument!= None and len(otpDocument) > 0:
        otpDocument['used'] = True
        db.otps.update_one({'_id':otpDocument['_id']}, {"$set": otpDocument}, upsert=False)
        file1 = open('docsearch/files/'+unit+'/index.txt', 'r')
        Lines = file1.readlines()
        idf = {}
        docTf = {}
        for line in Lines:
            lineArray = line.split("???")
            u=1
            idf[lineArray[0]] = float(lineArray[len(lineArray)-1])
            tf = {}
            for i in range(0,int(len(lineArray)/2)-1):
                tf[lineArray[0]] = int(lineArray[u+1])
                if lineArray[u] not in docTf.keys():
                    docTf[lineArray[u]] = tf
                else:
                    tf.update(docTf[lineArray[u]])
                    docTf[lineArray[u]] = tf
                u = u+2
                tf = {}
            u = 0

        result = removeStopWordsAndStemQuery(word_tokenize(query.encode('utf-16', 'surrogatepass').decode('utf-16')))

        queryVector = {}
        score = {}
        for key in idf:
            if key in result[1]:
                queryVector[key] = idf[key]*result[1][key]

        docVector = {}

        for key in docTf:
            tempTf = docTf[key]
            for token in tempTf:
                score[token] = tempTf[token]*idf[token]
            docVector[key] = score
            tempTf = {}
            score = {}
        simScore = 0.0
        for key in docVector:
            for key1 in queryVector:
                vector = docVector[key]
                if key1 in vector:
                    simScore = simScore + (vector[key1]*queryVector[key1])
            divisor = normalizedVectorValues(docVector[key], queryVector)
            if divisor > 0:
                simScore = simScore/divisor
            score[key] = simScore
            simScore = 0.0

        result = {}
        a1_sorted_keys = sorted(score, key=score.get, reverse=True)
        for r in a1_sorted_keys:
            result[r] = score[r]
        return(json.dumps(result))
    else:
        return '',403


@app.route('/document/<unit>/<otp>', methods=['GET'])
def getDocument(unit,otp):
    otpDocument = db.otps.find_one({'otpString': otp, 'used': False, 'role': 'Admin'})
    if otpDocument!= None and len(otpDocument) > 0:
        otpDocument['used'] = True
        db.otps.update_one({'_id':otpDocument['_id']}, {"$set": otpDocument}, upsert=False)
        pdfFileLocations = glob.glob("docsearch/files/"+unit+"/PDF/*.pdf")
        pptFileLocations = glob.glob("docsearch/files/"+unit+"/ppt/*.pptx")
        return({ "files" : pdfFileLocations + pptFileLocations})
    else:
        return '',403

@app.route('/document/upload/<unit>/<otp>', methods=['POST'])
def uploadDocuments(unit,otp):
    otpDocument = db.otps.find_one({'otpString': otp, 'used': False, 'role': 'Admin'})
    if otpDocument!= None and len(otpDocument) > 0:
        otpDocument['used'] = True
        db.otps.update_one({'_id':otpDocument['_id']}, {"$set": otpDocument}, upsert=False)
        indexDoc = db.index.find_one({'unit': unit, 'status': 'In Progress'})
        if indexDoc!= None and len(indexDoc) > 0:
            return 'In Progress', 405

        if 'file' not in request.files:
            print('No file part')
            return 'Error', 500

        files = request.files.getlist('file')
        for f in files:
            if f.filename == '':
                print('No selected file')
                return "error"
            if f and allowed_file(f.filename):
                filename = secure_filename(f.filename)
                if filename.rsplit('.', 1)[1].lower() == "pdf":
                    if not os.path.exists("docsearch/files/"+unit+"/PDF"):
                        os.makedirs("docsearch/files/"+unit+"/PDF")
                    f.save(os.path.join("docsearch/files/"+unit+"/PDF", filename))
                else:
                    if not os.path.exists("docsearch/files/"+unit+"/ppt"):
                        os.makedirs("docsearch/files/"+unit+"/ppt")
                    f.save(os.path.join("docsearch/files/"+unit+"/ppt", filename))
        return 'success', 200
    else:
        return '',403


@app.route('/document/download/<unit>/<name>/<otp>', methods=['GET'])
def downloadDocument(unit,name,otp):
    otpDocument = db.otps.find_one({'otpString': otp, 'used': False}) 
    if  otpDocument!= None and len(otpDocument) > 0:
        otpDocument['used'] = True
        db.otps.update_one({'_id':otpDocument['_id']}, {"$set": otpDocument}, upsert=False)
        if name.rsplit('.', 1)[1].lower() == "pdf":
            return send_file("docsearch/files/"+unit+"/PDF/"+name, as_attachment=True)
        else:
            return send_file("docsearch/files/"+unit+"/ppt/"+name, as_attachment=True)
    else:
        return '',403

@app.route('/document/delete/<otp>', methods=['DELETE'])
def deleteDocument(otp):
    otpDocument = db.otps.find_one({'otpString': otp, 'used': False, 'role': 'Admin'})
    if  otpDocument!= None and len(otpDocument) > 0:
        otpDocument['used'] = True
        db.otps.update_one({'_id':otpDocument['_id']}, {"$set": otpDocument}, upsert=False)
        indexDoc = db.index.find_one({'unit': request.json['unit'],'status': 'In Progress'})
        if indexDoc!= None and len(indexDoc) > 0:
            return("Indexing in Progress")
        if request.json['documentName'].rsplit('.', 1)[1].lower() == "pdf":
            try:
                os.remove("docsearch/files/"+request.json['unit']+"/PDF/"+request.json['documentName'])
                indexDocs(request.json['unit'])
                return "success"
            except Exception as e:
                return "failed"
        else:
            try:
                os.remove("docsearch/files/"+request.json['unit']+"/ppt/"+request.json['documentName'])
                return "success"
            except Exception as e:
                return "failed"
    else:
        return '',403

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0')